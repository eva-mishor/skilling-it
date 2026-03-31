"""Generate PPTX slides for 'Writing Claude Code Skills' webinar."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # Deep navy
BG_MEDIUM = RGBColor(0x22, 0x22, 0x3A)      # Slightly lighter
ACCENT = RGBColor(0xD4, 0x8A, 0x3E)         # Warm gold
ACCENT_LIGHT = RGBColor(0xE8, 0xB0, 0x6A)   # Light gold
TEXT_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
TEXT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
TEXT_DIM = RGBColor(0x88, 0x88, 0xAA)
CODE_BG = RGBColor(0x16, 0x16, 0x28)
GREEN = RGBColor(0x6B, 0xC9, 0x6B)
RED_SOFT = RGBColor(0xE0, 0x6C, 0x6C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=TEXT_WHITE, bullet_color=ACCENT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_code_block(slide, left, top, width, height, code, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_LIGHT
        p.font.name = "Courier New"
        p.space_after = Pt(2)
    return tf


def add_table_slide(slide, left, top, width, rows_data, col_widths=None,
                    font_size=15, header_color=ACCENT):
    cols = len(rows_data[0])
    rows = len(rows_data)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(rows * 0.45))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = BG_DARK
                else:
                    paragraph.font.color.rgb = TEXT_WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_MEDIUM if r % 2 == 1 else RGBColor(0x1E, 0x1E, 0x34)
    return table


def add_accent_line(slide, top):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(2), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_slide_number(slide, num, total=10):
    add_text_box(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=12, color=TEXT_DIM,
                 alignment=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
             "Writing Claude Code Skills", font_size=48, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(3.0))
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(1),
             "What they are, when you need one, and how to write one well",
             font_size=24, color=TEXT_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             "10-Minute Webinar", font_size=16, color=TEXT_DIM)

add_speaker_notes(slide,
    "Welcome. This is a 10-minute walkthrough of Claude Code skills — what they are, "
    "when you need one, how to structure them well, and how to share them with your team.\n\n"
    "By the end you'll know enough to build your first skill today.\n\n"
    "Target audience: anyone already using Claude Code who wants to go beyond one-off prompts."
)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 1)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Problem", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1.2),
             "Every new session, you explain your process again. From scratch.",
             font_size=26, color=TEXT_WHITE, bold=True)

# Two columns
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(0.5),
             "CLAUDE.md", font_size=22, color=TEXT_GRAY, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(3.5), Inches(5.2), Inches(3),
                ["Persistent project context",
                 "Always loaded, applies to everything",
                 "Passive \u2014 no structure, no phases",
                 "Can't say: 'follow this process'"],
                font_size=18, color=TEXT_GRAY)

add_text_box(slide, Inches(6.8), Inches(3.0), Inches(5.5), Inches(0.5),
             "Skills", font_size=22, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(3.5), Inches(5.5), Inches(3),
                ["Structured, multi-step processes",
                 "Loaded on demand, when triggered",
                 "Own files, tools, templates, logic",
                 "Decision points, gates, deliverables"],
                font_size=18, color=TEXT_WHITE)

add_speaker_notes(slide,
    "You're using Claude Code every day. You've got a workflow — maybe PR reviews, "
    "debugging production issues, onboarding new hires to a codebase.\n\n"
    "Every new session, you explain it again from scratch.\n\n"
    "CLAUDE.md helps — it gives Claude persistent project context. But it's passive. "
    "Always loaded, applies to everything, no structure. It can't say: 'when the user "
    "asks for X, follow this specific multi-step process, read these reference files at "
    "the right moments, produce these deliverables, and enforce these quality gates.'\n\n"
    "That's the gap skills fill.\n\n"
    "Think of it this way: CLAUDE.md is your team's coding standards. "
    "A skill is your team's playbook for a specific operation."
)

# ============================================================
# SLIDE 3: Why Skills?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 2)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Why Skills?", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
             "Five reasons to invest in skills",
             font_size=22, color=TEXT_GRAY)

# Five reasons — compact vertical list with number + title + description
reasons = [
    ("1", "Impose Determinism",
     "LLMs are non-deterministic. Skills give you structure \u2014 "
     "defined steps, gates, and quality bars that hold regardless of how the model feels today."),
    ("2", "Give the Right Guidelines",
     "Claude knows how to code. It doesn't know YOUR process, YOUR gotchas, YOUR conventions. "
     "Skills inject the specific knowledge that pushes it past default behavior."),
    ("3", "Centralize Learning",
     "Every mistake becomes a gotcha. Every session makes the skill better. "
     "Knowledge compounds across your team, not trapped in one conversation."),
    ("4", "Scale Across the Team",
     "One person figures out the process, everyone gets it. "
     "New hire day 1 runs the same skill as a 5-year veteran. Commit to the repo \u2014 the whole team is aligned."),
    ("5", "Context Efficiency",
     "Stop burning tokens re-explaining your process every session. "
     "Progressive disclosure loads only what's needed \u2014 Phase 1 doesn't pay for Phase 6."),
]

y = Inches(2.3)
for num, title, desc in reasons:
    add_text_box(slide, Inches(0.8), y, Inches(0.5), Inches(0.4),
                 num, font_size=24, color=ACCENT, bold=True)
    add_text_box(slide, Inches(1.4), y, Inches(3.5), Inches(0.4),
                 title, font_size=20, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(5.2), y, Inches(7.5), Inches(0.8),
                 desc, font_size=15, color=TEXT_GRAY)
    y += Inches(1.0)

add_speaker_notes(slide,
    "Before we get into how — the why.\n\n"
    "FIVE REASONS TO INVEST IN SKILLS:\n\n"
    "1. IMPOSE DETERMINISM IN A NON-DETERMINISTIC WORLD\n"
    "LLMs are stochastic. Same prompt, different day, different result. Skills are your "
    "way to impose structure: defined steps, gates, quality bars, deliverable formats. "
    "The model can be creative within the constraints, but the process itself is "
    "deterministic. Phase 1 always happens before Phase 2. The gate always gets checked. "
    "The deliverable always matches the template.\n\n"
    "2. GIVE THE RIGHT GUIDELINES\n"
    "Claude knows how to code. It doesn't know YOUR process, YOUR domain gotchas, YOUR "
    "org's conventions, YOUR customers' edge cases. Generic programming advice is wasted "
    "tokens. Skills inject the specific knowledge that pushes Claude past its defaults — "
    "the things it consistently gets wrong without guidance. A 'frontend design' skill "
    "doesn't teach CSS — it teaches YOUR team's taste, like avoiding Inter fonts.\n\n"
    "3. CENTRALIZE LEARNING — IMPROVE OVER TIME\n"
    "This is the compounding effect. Every mistake Claude makes becomes a gotcha. Every "
    "session makes the skill better. And because the skill lives in a shared location "
    "(repo, personal directory, enterprise), the learning compounds across your entire "
    "team — not trapped in one person's CLAUDE.md or one conversation's context window.\n\n"
    "Without skills, every team member re-discovers the same failure modes independently. "
    "With skills, someone hits the failure once, adds the gotcha, and nobody hits it again.\n\n"
    "4. SCALE ACROSS THE TEAM\n"
    "This is the distribution argument. One person figures out the process — the best way "
    "to review PRs, debug a specific service, onboard to a codebase. They codify it as a "
    "skill. Now everyone gets it. A new hire on day 1 runs the same skill as a 5-year "
    "veteran. Commit the skill to the repo and the whole team is aligned — same process, "
    "same quality bar, no tribal knowledge.\n\n"
    "3 and 4 are related but distinct: #3 is about accumulation (the skill gets better), "
    "#4 is about distribution (more people benefit). Together they're a flywheel — "
    "more users means more gotchas, which means a better skill, which means more users.\n\n"
    "5. CONTEXT EFFICIENCY\n"
    "Every time you explain your process in a new session, you're burning tokens on "
    "instructions instead of work. Skills eliminate that overhead. And progressive disclosure "
    "means you're not even loading the full skill — just the parts relevant to the current "
    "step. A 6-phase coaching process doesn't load Phase 6 frameworks while you're still in "
    "Phase 1. This is practical: less context spent on instructions = more context available "
    "for the actual work."
)

# ============================================================
# SLIDE 4: What Is a Skill?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 3)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "What Is a Skill?", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
             "A skill is a folder \u2014 not a prompt.",
             font_size=26, color=TEXT_WHITE, bold=True)

add_code_block(slide, Inches(0.8), Inches(2.5), Inches(6), Inches(3.5),
               ".claude/skills/my-skill/\n"
               "\u251c\u2500\u2500 SKILL.md              # Core instructions\n"
               "\u251c\u2500\u2500 references/            # Deep docs (on demand)\n"
               "\u2502   \u2514\u2500\u2500 framework.md\n"
               "\u251c\u2500\u2500 assets/                # Templates\n"
               "\u2502   \u2514\u2500\u2500 report-template.md\n"
               "\u2514\u2500\u2500 scripts/               # Executable tools\n"
               "    \u2514\u2500\u2500 validate.sh",
               font_size=16)

# Right side explanation
add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5.2), Inches(0.5),
             "Progressive disclosure:", font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(7.3), Inches(3.1), Inches(5.5), Inches(3),
                ["SKILL.md loads on trigger (~200 lines)",
                 "references/ loaded at specific steps",
                 "assets/ used when producing output",
                 "scripts/ run when Claude needs tools",
                 "",
                 "Phase 1 doesn't burn tokens on\nPhase 6 content"],
                font_size=17, color=TEXT_GRAY)

add_speaker_notes(slide,
    "A skill is a folder, not a prompt. At minimum it contains a SKILL.md with YAML "
    "frontmatter and markdown instructions. But the power is the folder structure.\n\n"
    "SKILL.md is the entry point — process flow, decision logic, quality gates. "
    "It tells Claude: here's the process, here's when to read deeper files.\n\n"
    "references/ is where detailed knowledge lives — framework tables, scoring rubrics, "
    "API docs. Claude reads these ONLY when it reaches the step that needs them. "
    "This is progressive disclosure — you're not burning context on content that's "
    "irrelevant to the current step.\n\n"
    "assets/ holds deliverable templates. scripts/ holds executable tools Claude can run.\n\n"
    "[DEMO: Show a real skill folder in the terminal with 'tree .claude/skills/']"
)

# ============================================================
# SLIDE 5: When Do You Need a Skill?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 4)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "When Do You Need a Skill?", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

# Two columns: need vs don't need
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
             "You need a skill when:", font_size=22, color=GREEN, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5),
                ["\u2713  Same multi-step process, across sessions",
                 "\u2713  Decision points, gates, conditional logic",
                 "\u2713  Reference material too large for CLAUDE.md",
                 "\u2713  Others should run the same process",
                 "\u2713  Specific deliverables in specific formats"],
                font_size=17, color=TEXT_WHITE)

add_text_box(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "You don't need a skill when:", font_size=22, color=RED_SOFT, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(2.5),
                ["\u2717  A CLAUDE.md instruction covers it",
                 "\u2717  It's a one-off task",
                 "\u2717  It's general knowledge Claude has"],
                font_size=17, color=TEXT_GRAY)

# Skill types table
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
             "Nine skill types:", font_size=20, color=ACCENT, bold=True)
add_table_slide(slide, Inches(0.8), Inches(5.0), Inches(11.5),
                [["Type", "Example"],
                 ["Business Process", "Multi-step workflows (release process, standups, ticket creation)"],
                 ["Data Fetching & Analysis", "Connect to data/monitoring stacks, produce reports"],
                 ["Code Quality & Review", "Enforce standards, adversarial review, style enforcement"],
                 ["Library/API Reference", "Correct usage of internal libraries, CLIs, design systems"],
                 ["Scaffolding & Templates", "Generate boilerplate (migrations, new apps, workflows)"],
                 ["Runbook", "Symptom \u2192 investigation \u2192 structured report"],
                 ["Product Verification", "Test/verify functionality (Playwright, video recording)"],
                 ["CI/CD & Deployment", "PR babysitting, gradual rollouts, cherry-pick workflows"],
                 ["Infrastructure Ops", "Resource cleanup, dependency mgmt, cost investigation"]],
                col_widths=[Inches(3.2), Inches(8.3)],
                font_size=12)

add_speaker_notes(slide,
    "Not everything needs to be a skill. Here's the decision filter.\n\n"
    "LEFT COLUMN — you need a skill when you repeat a multi-step process across sessions, "
    "when there are decision points or gates, when reference material is too large for "
    "CLAUDE.md, when other people need to run the same process, or when you want specific "
    "deliverables in specific formats.\n\n"
    "RIGHT COLUMN — you don't need a skill when a CLAUDE.md instruction covers it "
    "(like 'always use conventional commits'), when it's a one-off, or when it's general "
    "knowledge Claude already has.\n\n"
    "BOTTOM TABLE — nine skill types (from Anthropic's internal taxonomy). "
    "The important thing: pick ONE. Skills that straddle two types are confused.\n\n"
    "The three types people miss: Product Verification (pair with Playwright to have Claude "
    "record videos of output and include programmatic assertions), CI/CD & Deployment "
    "(PR babysitting, gradual rollouts — skills can reference other skills), and "
    "Infrastructure Ops (resource cleanup with guardrails for destructive actions).\n\n"
    "Business Process is the most common. Runbooks are underrated — "
    "take a Slack thread or alert as input, multi-tool investigation, structured report.\n\n"
    "[ASK AUDIENCE: 'What process do you repeat most often with Claude?' — that's their first skill.]"
)

# ============================================================
# SLIDE 6: The SKILL.md File
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 5)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The SKILL.md File", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

# Left: frontmatter code
add_code_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5),
               "---\n"
               "name: deploy-checklist\n"
               "description: \"Run pre-deploy validation\n"
               "  for production releases. Triggers on\n"
               "  'deploy to prod', 'release checklist'.\"\n"
               "allowed-tools: Bash, Read, Grep\n"
               "---",
               font_size=15)

# Right: three rules
add_text_box(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Three things to get right:", font_size=20, color=ACCENT, bold=True)

rules = [
    ("1. name = your slash command", "/deploy-checklist. Kebab-case, lowercase."),
    ("2. description is for Claude", "Not a summary \u2014 a trigger specification.\nInclude trigger phrases. Front-load the use case."),
    ("3. allowed-tools = permissions", "Restrict what Claude can do without asking.\nRead-only skill? Don't give write access."),
]
y = Inches(2.2)
for title, desc in rules:
    add_text_box(slide, Inches(6.8), y, Inches(5.5), Inches(0.35),
                 title, font_size=17, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(6.8), y + Inches(0.35), Inches(5.5), Inches(0.6),
                 desc, font_size=14, color=TEXT_GRAY)
    y += Inches(1.1)

# Bottom: other frontmatter
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.4),
             "Other useful frontmatter:", font_size=18, color=TEXT_GRAY, bold=True)
add_code_block(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2),
               "disable-model-invocation: true    # Only manual /slash-command\n"
               "user-invocable: false             # Only Claude can invoke (helper)\n"
               "context: fork                     # Run in isolated subagent\n"
               "paths: \"src/**/*.ts\"              # Only activate for matching files",
               font_size=14)

add_speaker_notes(slide,
    "The SKILL.md file has two parts: YAML frontmatter and a markdown body.\n\n"
    "FRONTMATTER — three things to get right:\n\n"
    "1. 'name' becomes your slash command. Whatever you put here — like 'deploy-checklist' — "
    "becomes /deploy-checklist. Kebab-case, lowercase.\n\n"
    "2. 'description' is for Claude, not for humans. This is critical. Claude scans every "
    "skill description at session start to decide which one matches your request. "
    "This is a trigger specification, not a summary.\n"
    "   Bad: 'A tool for deployments'\n"
    "   Good: 'Run pre-deploy validation for production releases. Triggers on deploy to prod, "
    "release checklist, pre-deploy check.'\n"
    "Front-load the use case. Include trigger phrases.\n\n"
    "3. 'allowed-tools' controls what Claude can do without asking permission. "
    "If your skill only reads files, restrict it — don't give write access to an audit skill.\n\n"
    "BOTTOM — other useful frontmatter:\n"
    "- disable-model-invocation: only manual /slash-command, Claude can't auto-trigger\n"
    "- user-invocable: false makes it a helper skill only Claude can call\n"
    "- context: fork runs the skill in an isolated subagent\n"
    "- paths: only activate when editing matching files\n\n"
    "THE BODY — most important rule: don't state the obvious. Claude already knows how to "
    "code. Focus on what pushes it out of default behavior — domain gotchas, org conventions, "
    "things it consistently gets wrong.\n\n"
    "[DEMO: Show frontmatter, then trigger the skill with a natural language request, "
    "show it activating automatically.]"
)

# ============================================================
# SLIDE 7: Progressive Disclosure
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 6)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Progressive Disclosure", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "The single most important structural decision in skill design",
             font_size=20, color=TEXT_GRAY)

# Anti-pattern
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.8),
             "A 2,000-line SKILL.md is not a skill. It's a context bomb.",
             font_size=24, color=RED_SOFT, bold=True)

# Table: what goes where
add_table_slide(slide, Inches(0.8), Inches(3.3), Inches(11.5),
                [["Location", "Content", "When Loaded"],
                 ["SKILL.md", "Process flow, logic, gates, gotchas", "Always (on trigger)"],
                 ["references/", "Detailed frameworks, data, rubrics", "On demand (at specific steps)"],
                 ["assets/", "Deliverable templates", "When producing output"],
                 ["scripts/", "Executable tools", "When Claude needs to run them"]],
                col_widths=[Inches(2.5), Inches(5.5), Inches(3.5)],
                font_size=15)

add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
             "SKILL.md is the table of contents and decision logic.\nThe chapters live in references/.",
             font_size=20, color=ACCENT)

add_speaker_notes(slide,
    "This is the single most important structural decision in skill design.\n\n"
    "A 2,000-line SKILL.md is a context bomb. Claude loads the entire thing when the skill "
    "triggers. Every token is present whether it's relevant to the current step or not.\n\n"
    "The pattern: SKILL.md stays compact — process flow, decision points, gate criteria, "
    "gotchas, and a directory listing of what's available. It says things like:\n"
    "   'Read references/scoring-framework.md for the full threshold table.'\n"
    "   'Use the template in assets/report-template.md for the deliverable.'\n\n"
    "Claude reads those files only when it reaches that step.\n\n"
    "TABLE — what goes where:\n"
    "- SKILL.md: process, logic, gates, gotchas — always loaded on trigger\n"
    "- references/: detailed frameworks, data, rubrics — loaded on demand at specific steps\n"
    "- assets/: deliverable templates — loaded when producing output\n"
    "- scripts/: executable tools — loaded when Claude needs to run them\n\n"
    "The metaphor: SKILL.md is the table of contents and decision logic. "
    "The chapters live in references/.\n\n"
    "When extracting from a monolithic SKILL.md: preserve exact phrases and section headings — "
    "renaming during extraction loses content traceability. Reference files must stand alone — "
    "add a one-line intro explaining context.\n\n"
    "[DEMO: Show Claude reading a reference file mid-process — the moment where it says "
    "'Let me read references/...' in a transcript.]"
)

# ============================================================
# SLIDE 8: Five Writing Rules
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 7)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Five Rules That Matter", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

rules_data = [
    ("1. Build a gotchas section",
     "Highest-signal content = what goes wrong. Every time Claude\n"
     "makes a mistake, add it. After 10 sessions, this is your\n"
     "most valuable section."),
    ("2. Avoid railroading",
     "Principles and constraints, not rigid scripts. Rigid instructions\n"
     "break the moment the situation doesn't match."),
    ("3. Do the work first, then red/green/refactor",
     "Live session \u2192 first draft \u2192 Red (find gaps) \u2192 Green (fix) \u2192\n"
     "Refactor \u2192 Data audit \u2192 Ship. Use /skill-creator to scaffold."),
    ("4. Audit every factual claim",
     "Grep for numbers and assertions. Search for sources.\n"
     "If unsourced: reframe as methodology or cut it."),
    ("5. One type, one skill",
     "If it straddles two types, it's confused about what it does.\n"
     "Split it."),
]

y = Inches(1.6)
for title, desc in rules_data:
    add_text_box(slide, Inches(0.8), y, Inches(11), Inches(0.35),
                 title, font_size=20, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.38), Inches(10.5), Inches(0.75),
                 desc, font_size=15, color=TEXT_GRAY)
    y += Inches(1.1)

add_speaker_notes(slide,
    "Five rules that separate skills that work from skills that don't.\n\n"
    "1. BUILD A GOTCHAS SECTION\n"
    "The highest-signal content in any skill is what goes wrong. Every time Claude makes "
    "a mistake while using the skill, add it as a gotcha. The example table shows real "
    "gotchas: asking about pricing before pain, citing training data as research, skipping "
    "gates. After 10 sessions, this section is the most valuable part of your skill. "
    "It's a living document — it grows from real failures.\n\n"
    "2. AVOID RAILROADING\n"
    "Give Claude principles and constraints, not rigid scripts. Rigid step-by-step "
    "instructions break the moment the situation doesn't match — and it never matches "
    "exactly. Be especially careful with reusable skills: rigidity that works for one "
    "case fails for the next.\n\n"
    "3. DO THE WORK FIRST, THEN RED/GREEN/REFACTOR\n"
    "Never write a skill abstractly. The creation process:\n"
    "   Live session (real case, real decisions)\n"
    "   -> First draft (extract what worked)\n"
    "   -> Red (what's wrong? test on a second case, watch it fail)\n"
    "   -> Green (fix the gaps, add gotchas)\n"
    "   -> Refactor (clean up, extract references, align structure)\n"
    "   -> Data audit (verify every factual claim)\n"
    "   -> Ship\n"
    "You can use /skill-creator to scaffold the initial structure, then refine through "
    "this cycle. One of my skills missed 5 structural issues — no rescue loop, no kill "
    "hierarchy, no onboarding flow — only visible when tested on a real case.\n\n"
    "4. AUDIT EVERY FACTUAL CLAIM\n"
    "After writing, grep for specific numbers, example classifications, and 'X is Y' "
    "assertions. Search for sources. If unsourced, find data or reframe as methodology. "
    "One first draft had 8 unsourced claims presented as facts — wrong taxonomy examples, "
    "a made-up '80%' stat. All caught before shipping.\n\n"
    "5. ONE TYPE, ONE SKILL\n"
    "If it's trying to be a code reviewer AND a scaffolding tool, split it. Skills that "
    "straddle types are confused about what they do and do neither well.\n\n"
    "[DEMO: Live /skill-creator — describe a process, show the scaffold, walk through "
    "how you'd refine it through red/green/refactor.]"
)

# ============================================================
# SLIDE 9: Sharing & Installing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 8)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Sharing & Installing", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

# Scope table
add_table_slide(slide, Inches(0.8), Inches(1.6), Inches(11.5),
                [["Scope", "Path", "Who Gets It"],
                 ["Personal", "~/.claude/skills/<name>/", "Just you, all projects"],
                 ["Project", ".claude/skills/<name>/", "Everyone on this repo"],
                 ["Enterprise", "Managed settings", "All org users"]],
                col_widths=[Inches(2.5), Inches(5), Inches(4)],
                font_size=16)

# GitHub sharing
add_text_box(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.5),
             "Share via GitHub:", font_size=22, color=ACCENT, bold=True)
add_code_block(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.8),
               "claude install-skill https://github.com/you/repo/tree/main/skills/my-skill",
               font_size=16)

# Project skills
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.5),
             "Project skills \u2014 commit to repo:", font_size=22, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1),
             "Commit .claude/skills/ to your repo. Every teammate who clones gets\n"
             "the skills automatically. No install step. This is how you standardize\n"
             "processes across a team.",
             font_size=18, color=TEXT_GRAY)

add_speaker_notes(slide,
    "Skills live in one of three places, and location determines scope.\n\n"
    "TABLE:\n"
    "- Personal (~/.claude/skills/): just you, available in all your projects\n"
    "- Project (.claude/skills/): everyone who clones the repo gets them automatically\n"
    "- Enterprise (managed settings): all org users\n\n"
    "SHARING VIA GITHUB\n"
    "Put your skill folder in a public repo. Anyone installs with one command:\n"
    "   claude install-skill https://github.com/you/repo/tree/main/skills/my-skill\n"
    "It clones into their personal skills directory. Available in every session.\n\n"
    "PROJECT SKILLS — THE TEAM PLAY\n"
    "This is the most powerful pattern for teams. Commit .claude/skills/ to your repo. "
    "Every teammate who clones gets the skills automatically — no install step. "
    "This is how you standardize processes: your deploy checklist, PR review process, "
    "incident response runbook. Same process, every person, every time.\n\n"
    "ADVANCED: INTERNAL MARKETPLACE\n"
    "For larger teams, Anthropic recommends a marketplace pattern: find useful skills "
    "organically, upload sandbox versions to GitHub, move to an internal marketplace via PR "
    "after gaining traction. Curate before release to prevent bad or redundant skills.\n\n"
    "SKILL COMPOSITION\n"
    "Skills can reference other skills by name. If the referenced skill is installed, "
    "Claude will invoke it. Think of it as dependency management — a deploy skill can "
    "call a PR-review skill before pushing.\n\n"
    "MEASURING SUCCESS\n"
    "Use a PreToolUse hook to log when skills trigger. Track popular vs. under-triggering "
    "skills. If a skill never triggers automatically, its description needs work.\n\n"
    "Skills also follow the Agent Skills open standard (agentskills.io), so the format "
    "works with compatible tools beyond Claude Code — though frontmatter fields vary.\n\n"
    "[DEMO: Live 'claude install-skill' from a GitHub URL — show it arriving in "
    "~/.claude/skills/ and being available immediately.]"
)

# ============================================================
# SLIDE 10: Recap + CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 9)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Cheat Sheet", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

recap_items = [
    ("A skill is a folder", "SKILL.md + references + assets + scripts"),
    ("Progressive disclosure", "SKILL.md is the router; references are the depth"),
    ("Description triggers the skill", "Write it for Claude, not for humans"),
    ("Gotchas are your best content", "Build them from real failures over time"),
    ("Do the work first", "Live session \u2192 red/green/refactor \u2192 ship"),
]

y = Inches(1.6)
for title, desc in recap_items:
    add_text_box(slide, Inches(0.8), y, Inches(5), Inches(0.35),
                 title, font_size=20, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(5.5), y, Inches(7), Inches(0.35),
                 desc, font_size=18, color=TEXT_GRAY)
    y += Inches(0.55)

# CTA box
cta_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5)
)
cta_shape.fill.solid()
cta_shape.fill.fore_color.rgb = BG_MEDIUM
cta_shape.line.color.rgb = ACCENT
cta_shape.line.width = Pt(2)

tf = cta_shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.4)
tf.margin_top = Inches(0.3)

p = tf.paragraphs[0]
p.text = "Start here:"
p.font.size = Pt(22)
p.font.color.rgb = ACCENT
p.font.bold = True
p.font.name = "Calibri"

p2 = tf.add_paragraph()
p2.text = ""
p2.space_after = Pt(4)

steps = [
    "1. Open Claude Code, type /skill-creator",
    "2. Describe the process you repeat most often",
    "3. Use the generated skill tomorrow",
    "4. Add a gotcha when it gets something wrong",
    "5. Iterate"
]
for step in steps:
    p3 = tf.add_paragraph()
    p3.text = step
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_WHITE
    p3.font.name = "Calibri"
    p3.space_after = Pt(4)

add_speaker_notes(slide,
    "Cheat sheet — five things to remember:\n\n"
    "1. A skill is a folder — SKILL.md + references + assets + scripts\n"
    "2. Progressive disclosure — SKILL.md is the router, references are the depth\n"
    "3. The description triggers the skill — write it for Claude, not humans\n"
    "4. Gotchas are your best content — build from real failures over time\n"
    "5. Do the work first — live session, then red/green/refactor, then ship\n\n"
    "CALL TO ACTION:\n"
    "Your first skill should be the process you're most tired of re-explaining. "
    "The one where you think 'I wish Claude just knew how we do this.'\n\n"
    "You don't have to start from scratch. Claude Code has /skill-creator — "
    "describe what you want, it generates the skeleton.\n\n"
    "The loop: /skill-creator -> use it tomorrow -> add a gotcha when it gets something "
    "wrong -> iterate. That's it.\n\n"
    "COMMON Q&A:\n\n"
    "Q: How is this different from system prompts?\n"
    "A: System prompts are static text. Skills are structured toolkits — folders, "
    "progressive disclosure, scripts, conditional logic, contextual triggering.\n\n"
    "Q: Can skills call other skills?\n"
    "A: Not directly, but a skill can reference another by name and Claude will invoke it.\n\n"
    "Q: What about context window limits?\n"
    "A: That's exactly why progressive disclosure matters. A well-structured skill loads "
    "~200 lines initially and reads deeper files only when needed.\n\n"
    "Q: How do I know if my skill is working well?\n"
    "A: Use it. The gotchas section is your feedback loop. If you're adding gotchas every "
    "session, the skill is learning. When you stop adding them, it's mature.\n\n"
    "[ASK AUDIENCE: 'How many of you have a CLAUDE.md over 200 lines?' — "
    "some of that should probably be skills.]"
)

# ============================================================
# SLIDE 11: Resources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_slide_number(slide, 10)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Resources", font_size=36, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.2))

resources = [
    ("Anthropic \u2014 How We Use Skills (Thariq Shihipar)",
     "linkedin.com/pulse/lessons-from-building-claude-code-\nhow-we-use-skills-thariq-shihipar-iclmc"),
    ("Claude Code Documentation \u2014 Skills",
     "docs.anthropic.com/en/docs/claude-code/skills"),
    ("Agent Skills Open Standard",
     "agentskills.io"),
    ("Claude Code Documentation \u2014 Hooks",
     "docs.anthropic.com/en/docs/claude-code/hooks"),
    ("/skill-creator",
     "Built into Claude Code \u2014 type /skill-creator to scaffold a new skill"),
]

y = Inches(1.8)
for title, url in resources:
    add_text_box(slide, Inches(0.8), y, Inches(11), Inches(0.35),
                 title, font_size=18, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.35), Inches(10.5), Inches(0.5),
                 url, font_size=15, color=ACCENT_LIGHT)
    y += Inches(1.0)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             "Links in the speaker notes \u2014 copy-paste friendly",
             font_size=14, color=TEXT_DIM)

add_speaker_notes(slide,
    "RESOURCES — full URLs for copy-paste:\n\n"
    "1. Anthropic's internal lessons on skills (Thariq Shihipar, Anthropic):\n"
    "   https://www.linkedin.com/pulse/lessons-from-building-claude-code-how-we-use-skills-thariq-shihipar-iclmc/\n"
    "   Best single resource on skill design. Covers 9 skill types, gotchas, progressive "
    "disclosure, hooks, marketplace patterns, and measuring success. Written by someone "
    "building skills at Anthropic.\n\n"
    "2. Claude Code Skills Documentation:\n"
    "   https://docs.anthropic.com/en/docs/claude-code/skills\n"
    "   Official docs — frontmatter reference, installation, scope, and configuration.\n\n"
    "3. Agent Skills Open Standard:\n"
    "   https://agentskills.io\n"
    "   The open standard skills are built on. Useful if you want skills portable across "
    "AI tools.\n\n"
    "4. Claude Code Hooks Documentation:\n"
    "   https://docs.anthropic.com/en/docs/claude-code/hooks\n"
    "   How to register on-demand hooks that activate only when a skill runs. Examples: "
    "'/careful' blocks rm -rf, DROP TABLE, force-push; '/freeze' restricts edits to "
    "specific directories.\n\n"
    "5. /skill-creator:\n"
    "   Built into Claude Code. Type /skill-creator and describe your process — it generates "
    "the folder structure, frontmatter, and initial SKILL.md content.\n\n"
    "Thank you! Questions?"
)


# Save
output_path = "/Users/user/Documents/Code/the-shipping-playbook/sessions/shipping-playbook/webinar-claude-code-skills.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
